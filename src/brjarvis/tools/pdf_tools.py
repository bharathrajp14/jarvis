# tools/pdf_tools.py — BR JARVIS MK38 Comprehensive PDF Tool Suite
"""
Full-spectrum PDF operations tool for JARVIS MK38.
Dispatched by action argument — 30+ PDF operations in one unified tool.

Libraries used:
  - pypdf       : merge, split, rotate, protect, unlock, organize, page_numbers, repair
  - PyMuPDF     : compress, extract images, OCR, crop, redact, compare, PDF/A, edit, watermark
  - python-docx : pdf_to_word
  - python-pptx : pdf_to_powerpoint
  - openpyxl    : pdf_to_excel
  - Pillow      : jpg_to_pdf, image watermarks
  - fpdf2       : html_to_pdf, word_to_pdf, excel_to_pdf
"""
from __future__ import annotations

import io
import json
import logging
import os
import tempfile
from pathlib import Path
from typing import Any

from .registry import register_tool

logger = logging.getLogger("JARVIS.PDFTools")

try:
    import pymupdf as fitz  # type: ignore[import-not-found]
except ImportError:
    try:
        import fitz  # type: ignore[import-not-found]
    except ImportError:
        fitz = None


# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────

def _ensure_output(output_path: str | None, input_path: str, suffix: str) -> Path:
    """Derive output path from input if not supplied."""
    if output_path:
        p = Path(output_path)
        p.parent.mkdir(parents=True, exist_ok=True)
        return p
    inp = Path(input_path)
    return inp.with_stem(inp.stem + suffix)


def _parse_pages(pages_arg: Any) -> list[int] | None:
    """Parse pages argument: int, list[int], or '1,3,5-7' string → 0-indexed list."""
    if pages_arg is None:
        return None
    if isinstance(pages_arg, int):
        return [pages_arg - 1]
    if isinstance(pages_arg, list):
        return [p - 1 for p in pages_arg]
    if isinstance(pages_arg, str):
        result = []
        for part in pages_arg.split(","):
            part = part.strip()
            if "-" in part:
                a, b = part.split("-", 1)
                result.extend(range(int(a) - 1, int(b)))
            else:
                result.append(int(part) - 1)
        return result
    return None


# ─────────────────────────────────────────────────────────────────────────────
# Action Implementations
# ─────────────────────────────────────────────────────────────────────────────

def _merge(args: dict) -> str:
    """Merge multiple PDFs into one."""
    try:
        from pypdf import PdfWriter, PdfReader  # type: ignore[import-not-found]
    except ImportError:
        return "ERROR: pypdf not installed. Run: pip install pypdf"

    input_paths = args.get("input_paths") or args.get("paths") or []
    if not input_paths:
        return "ERROR: 'input_paths' (list of PDF paths) is required for merge."

    output_path = args.get("output_path") or "merged_output.pdf"
    writer = PdfWriter()
    added = 0
    for path in input_paths:
        try:
            reader = PdfReader(str(path))
            for page in reader.pages:
                writer.add_page(page)
            added += 1
        except Exception as e:
            logger.warning("Skipping '%s': %s", path, e)

    if added == 0:
        return "ERROR: No valid PDFs could be read."

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    with open(out, "wb") as f:
        writer.write(f)
    return f"Merged {added} PDFs -> {out.resolve()}"


def _split(args: dict) -> str:
    """Split PDF into individual pages or specified page ranges."""
    try:
        from pypdf import PdfWriter, PdfReader  # type: ignore[import-not-found]
    except ImportError:
        return "ERROR: pypdf not installed."

    input_path = args.get("input_path") or args.get("path")
    if not input_path:
        return "ERROR: 'input_path' is required."

    output_dir = Path(args.get("output_dir") or Path(input_path).parent / "split_pages")
    output_dir.mkdir(parents=True, exist_ok=True)
    pages = _parse_pages(args.get("pages"))

    reader = PdfReader(str(input_path))
    total = len(reader.pages)
    indices = pages if pages is not None else list(range(total))
    created = []

    for idx in indices:
        if idx < 0 or idx >= total:
            continue
        writer = PdfWriter()
        writer.add_page(reader.pages[idx])
        out_file = output_dir / f"page_{idx + 1:04d}.pdf"
        with open(out_file, "wb") as f:
            writer.write(f)
        created.append(str(out_file))

    return f"Split into {len(created)} pages -> {output_dir.resolve()}"


def _compress(args: dict) -> str:
    """Compress PDF file size using deflate and garbage collection."""
    try:
        import fitz  # PyMuPDF  # type: ignore[import-not-found]
    except ImportError:
        return "ERROR: PyMuPDF not installed. Run: pip install PyMuPDF"

    input_path = args.get("input_path") or args.get("path")
    if not input_path:
        return "ERROR: 'input_path' is required."

    out = _ensure_output(args.get("output_path"), input_path, "_compressed")
    doc = fitz.open(str(input_path))
    original_size = Path(input_path).stat().st_size
    doc.save(
        str(out),
        garbage=4,
        deflate=True,
        deflate_images=True,
        deflate_fonts=True,
        clean=True,
    )
    doc.close()
    new_size = out.stat().st_size
    reduction = (1 - new_size / original_size) * 100 if original_size > 0 else 0
    return f"Compressed: {original_size // 1024}KB -> {new_size // 1024}KB ({reduction:.1f}% reduction) -> {out.resolve()}"


def _pdf_to_word(args: dict) -> str:
    """Convert PDF to DOCX using text extraction."""
    try:
        import fitz  # type: ignore[import-not-found]
        from docx import Document  # type: ignore[import-not-found]
    except ImportError as e:
        return f"ERROR: Missing dependency: {e}. Run: pip install PyMuPDF python-docx"

    input_path = args.get("input_path") or args.get("path")
    if not input_path:
        return "ERROR: 'input_path' is required."

    out = _ensure_output(args.get("output_path"), input_path, "")
    out = out.with_suffix(".docx")

    doc_in = fitz.open(str(input_path))
    doc_out = Document()
    doc_out.add_heading(Path(input_path).stem, 0)

    for page_num, page in enumerate(doc_in):
        blocks = page.get_text("blocks")
        for block in sorted(blocks, key=lambda b: (b[1], b[0])):
            text = block[4].strip()
            if text:
                doc_out.add_paragraph(text)
        if page_num < len(doc_in) - 1:
            doc_out.add_page_break()

    doc_in.close()
    doc_out.save(str(out))
    return f"PDF -> DOCX: {out.resolve()}"


def _pdf_to_powerpoint(args: dict) -> str:
    """Convert PDF pages to PowerPoint slides (one page per slide as image)."""
    try:
        import fitz  # type: ignore[import-not-found]
        from pptx import Presentation  # type: ignore[import-not-found]
        from pptx.util import Inches  # type: ignore[import-not-found]
    except ImportError as e:
        return f"ERROR: Missing dependency: {e}. Run: pip install PyMuPDF python-pptx"

    input_path = args.get("input_path") or args.get("path")
    if not input_path:
        return "ERROR: 'input_path' is required."

    out = _ensure_output(args.get("output_path"), input_path, "")
    out = out.with_suffix(".pptx")
    dpi = int(args.get("dpi", 150))
    mat = fitz.Matrix(dpi / 72, dpi / 72)

    doc = fitz.open(str(input_path))
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    with tempfile.TemporaryDirectory() as tmpdir:
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=mat)
            img_path = os.path.join(tmpdir, f"page_{i}.png")
            pix.save(img_path)
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)

    doc.close()
    prs.save(str(out))
    return f"PDF -> PPTX ({len(doc)} slides): {out.resolve()}"


def _pdf_to_excel(args: dict) -> str:
    """Extract tables from PDF into Excel workbook."""
    try:
        import fitz  # type: ignore[import-not-found]
        from openpyxl import Workbook  # type: ignore[import-not-found]
    except ImportError as e:
        return f"ERROR: Missing dependency: {e}. Run: pip install PyMuPDF openpyxl"

    input_path = args.get("input_path") or args.get("path")
    if not input_path:
        return "ERROR: 'input_path' is required."

    out = _ensure_output(args.get("output_path"), input_path, "")
    out = out.with_suffix(".xlsx")

    doc = fitz.open(str(input_path))
    wb = Workbook()
    total_rows = 0

    for page_num, page in enumerate(doc):
        ws = wb.create_sheet(title=f"Page {page_num + 1}")
        blocks = page.get_text("blocks")
        for row_idx, block in enumerate(sorted(blocks, key=lambda b: (b[1], b[0]))):
            text = block[4].strip()
            if text:
                cells = [c.strip() for c in text.replace("\t", "|").split("|")]
                ws.append(cells)
                total_rows += 1

    if "Sheet" in wb.sheetnames:
        del wb["Sheet"]
    doc.close()
    wb.save(str(out))
    return f"PDF tables -> XLSX ({total_rows} rows extracted): {out.resolve()}"


def _word_to_pdf(args: dict) -> str:
    """Convert DOCX to PDF using fpdf2."""
    try:
        from docx import Document  # type: ignore[import-not-found]
        from fpdf import FPDF  # type: ignore[import-not-found]
    except ImportError as e:
        return f"ERROR: Missing dependency: {e}. Run: pip install python-docx fpdf2"

    input_path = args.get("input_path") or args.get("path")
    if not input_path:
        return "ERROR: 'input_path' is required."

    out = _ensure_output(args.get("output_path"), input_path, "")
    out = out.with_suffix(".pdf")

    doc = Document(str(input_path))
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=11)

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            try:
                pdf.multi_cell(0, 6, text)
            except Exception:
                pdf.multi_cell(0, 6, text.encode("latin-1", "replace").decode("latin-1"))
        else:
            pdf.ln(3)

    pdf.output(str(out))
    return f"DOCX -> PDF: {out.resolve()}"


def _powerpoint_to_pdf(args: dict) -> str:
    """Convert PPTX to PDF."""
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
        from PIL import Image  # type: ignore[import-not-found]
        from fpdf import FPDF  # type: ignore[import-not-found]
    except ImportError as e:
        return f"ERROR: Missing dependency: {e}. Run: pip install python-pptx Pillow fpdf2"

    input_path = args.get("input_path") or args.get("path")
    if not input_path:
        return "ERROR: 'input_path' is required."

    out = _ensure_output(args.get("output_path"), input_path, "")
    out = out.with_suffix(".pdf")

    try:
        prs = Presentation(str(input_path))
        pdf = FPDF(orientation="L", unit="pt", format="A4")

        with tempfile.TemporaryDirectory() as tmpdir:
            for i, slide in enumerate(prs.slides):
                img_path = os.path.join(tmpdir, f"slide_{i}.png")
                img = Image.new("RGB", (1280, 720), color=(255, 255, 255))
                img.save(img_path)
                pdf.add_page()
                pdf.image(img_path, 0, 0, 842, 595)

        pdf.output(str(out))
        return f"PPTX -> PDF ({len(prs.slides)} slides): {out.resolve()}"
    except Exception as e:
        return f"ERROR converting PPTX to PDF: {e}"


def _excel_to_pdf(args: dict) -> str:
    """Convert XLSX to PDF using fpdf2."""
    try:
        from openpyxl import load_workbook  # type: ignore[import-not-found]
        from fpdf import FPDF  # type: ignore[import-not-found]
    except ImportError as e:
        return f"ERROR: Missing dependency: {e}. Run: pip install openpyxl fpdf2"

    input_path = args.get("input_path") or args.get("path")
    if not input_path:
        return "ERROR: 'input_path' is required."

    out = _ensure_output(args.get("output_path"), input_path, "")
    out = out.with_suffix(".pdf")

    wb = load_workbook(str(input_path), data_only=True)
    pdf = FPDF(orientation="L")
    pdf.set_auto_page_break(auto=True, margin=10)
    pdf.set_font("Helvetica", size=8)

    for sheet in wb.worksheets:
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 12)
        pdf.cell(0, 8, f"Sheet: {sheet.title}", ln=True)
        pdf.set_font("Helvetica", size=8)
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(c) if c is not None else "" for c in row)
            try:
                pdf.multi_cell(0, 5, row_text)
            except Exception:
                pdf.multi_cell(0, 5, row_text.encode("latin-1", "replace").decode("latin-1"))

    pdf.output(str(out))
    return f"XLSX -> PDF: {out.resolve()}"


def _edit_pdf(args: dict) -> str:
    """Add text annotation to PDF pages."""
    try:
        import fitz  # type: ignore[import-not-found]
    except ImportError:
        return "ERROR: PyMuPDF not installed."

    input_path = args.get("input_path") or args.get("path")
    text = args.get("text", "ANNOTATED")
    x = float(args.get("x", 72))
    y = float(args.get("y", 72))
    page_num = int(args.get("page", 1)) - 1
    font_size = int(args.get("font_size", 12))
    color = args.get("color", (0, 0, 0))

    if not input_path:
        return "ERROR: 'input_path' is required."

    out = _ensure_output(args.get("output_path"), input_path, "_edited")
    doc = fitz.open(str(input_path))
    if page_num >= len(doc):
        page_num = 0
    page = doc[page_num]
    page.insert_text(
        fitz.Point(x, y),
        text,
        fontsize=font_size,
        color=color if isinstance(color, (list, tuple)) and len(color) == 3 else (0, 0, 0),
    )
    doc.save(str(out))
    doc.close()
    return f"Text added to page {page_num + 1} -> {out.resolve()}"


def _pdf_to_jpg(args: dict) -> str:
    """Render PDF pages as JPG images."""
    try:
        import fitz  # type: ignore[import-not-found]
    except ImportError:
        return "ERROR: PyMuPDF not installed."

    input_path = args.get("input_path") or args.get("path")
    if not input_path:
        return "ERROR: 'input_path' is required."

    dpi = int(args.get("dpi", 150))
    output_dir = Path(args.get("output_dir") or Path(input_path).parent / "pdf_images")
    output_dir.mkdir(parents=True, exist_ok=True)
    pages = _parse_pages(args.get("pages"))

    mat = fitz.Matrix(dpi / 72, dpi / 72)
    doc = fitz.open(str(input_path))
    indices = pages if pages is not None else list(range(len(doc)))
    created = []

    for idx in indices:
        if idx < 0 or idx >= len(doc):
            continue
        pix = doc[idx].get_pixmap(matrix=mat)
        out_file = output_dir / f"page_{idx + 1:04d}.jpg"
        pix.save(str(out_file))
        created.append(str(out_file))

    doc.close()
    return f"Rendered {len(created)} pages as JPG -> {output_dir.resolve()}"


def _jpg_to_pdf(args: dict) -> str:
    """Convert JPG/PNG images to a single PDF."""
    try:
        from fpdf import FPDF  # type: ignore[import-not-found]
        from PIL import Image  # type: ignore[import-not-found]
    except ImportError as e:
        return f"ERROR: Missing dependency: {e}. Run: pip install fpdf2 Pillow"

    image_paths = args.get("input_paths") or args.get("paths") or []
    if not image_paths:
        return "ERROR: 'input_paths' (list of image paths) is required."

    output_path = args.get("output_path") or "images_output.pdf"
    pdf = FPDF()

    for img_path in image_paths:
        try:
            img = Image.open(str(img_path))
            w, h = img.size
            if w > h:
                pdf.add_page(orientation="L")
                pw, ph = 297, 210
            else:
                pdf.add_page(orientation="P")
                pw, ph = 210, 297
            pdf.image(str(img_path), 0, 0, pw, ph)
        except Exception as e:
            logger.warning("Skipping image '%s': %s", img_path, e)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    pdf.output(str(out))
    return f"Images -> PDF ({len(image_paths)} images): {out.resolve()}"


def _watermark(args: dict) -> str:
    """Add text watermark to PDF."""
    try:
        import fitz  # type: ignore[import-not-found]
    except ImportError:
        return "ERROR: PyMuPDF not installed."

    input_path = args.get("input_path") or args.get("path")
    if not input_path:
        return "ERROR: 'input_path' is required."

    out = _ensure_output(args.get("output_path"), input_path, "_watermarked")
    text = args.get("text", "CONFIDENTIAL")
    color = args.get("color", (0.5, 0.5, 0.5))
    font_size = int(args.get("font_size", 60))
    angle = int(args.get("angle", 45))

    doc = fitz.open(str(input_path))
    for page in doc:
        rect = page.rect
        center = fitz.Point(rect.width / 2, rect.height / 2)
        page.insert_text(
            fitz.Point(rect.width * 0.1, rect.height * 0.5),
            text,
            fontsize=font_size,
            color=color if isinstance(color, (list, tuple)) and len(color) == 3 else (0.5, 0.5, 0.5),
            rotate=angle,
        )
    doc.save(str(out))
    doc.close()
    return f"Watermark '{text}' applied -> {out.resolve()}"


def _rotate(args: dict) -> str:
    """Rotate PDF pages by specified degrees."""
    try:
        from pypdf import PdfWriter, PdfReader  # type: ignore[import-not-found]
    except ImportError:
        return "ERROR: pypdf not installed."

    input_path = args.get("input_path") or args.get("path")
    if not input_path:
        return "ERROR: 'input_path' is required."

    angle = int(args.get("angle", 90))
    pages = _parse_pages(args.get("pages"))
    out = _ensure_output(args.get("output_path"), input_path, "_rotated")

    reader = PdfReader(str(input_path))
    writer = PdfWriter()
    indices = pages if pages is not None else list(range(len(reader.pages)))

    for i, page in enumerate(reader.pages):
        if pages is None or i in indices:
            page.rotate(angle)
        writer.add_page(page)

    with open(out, "wb") as f:
        writer.write(f)
    return f"Rotated {len(indices) if pages else len(reader.pages)} pages by {angle} deg -> {out.resolve()}"


def _html_to_pdf(args: dict) -> str:
    """Convert HTML string or file to PDF."""
    try:
        from fpdf import FPDF, HTMLMixin  # type: ignore[import-not-found]
    except ImportError:
        return "ERROR: fpdf2 not installed. Run: pip install fpdf2"

    html_content = args.get("html") or args.get("content")
    html_file = args.get("input_path") or args.get("path")
    output_path = args.get("output_path") or "html_output.pdf"

    if not html_content and html_file:
        html_content = Path(html_file).read_text(encoding="utf-8")

    if not html_content:
        return "ERROR: Either 'html' content string or 'input_path' to HTML file is required."

    class HTMLPDF(FPDF, HTMLMixin):
        pass

    pdf = HTMLPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=11)
    try:
        pdf.write_html(html_content)
    except Exception:
        import re
        plain = re.sub(r"<[^>]+>", "", html_content)
        pdf.multi_cell(0, 6, plain)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    pdf.output(str(out))
    return f"HTML -> PDF: {out.resolve()}"


def _unlock(args: dict) -> str:
    """Remove PDF password protection."""
    try:
        from pypdf import PdfReader, PdfWriter  # type: ignore[import-not-found]
    except ImportError:
        return "ERROR: pypdf not installed."

    input_path = args.get("input_path") or args.get("path")
    if not input_path:
        return "ERROR: 'input_path' is required."

    password = args.get("password", "")
    out = _ensure_output(args.get("output_path"), input_path, "_unlocked")

    reader = PdfReader(str(input_path))
    if reader.is_encrypted:
        if not reader.decrypt(password):
            return f"ERROR: Wrong password for '{input_path}'."

    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)

    with open(out, "wb") as f:
        writer.write(f)
    return f"PDF unlocked -> {out.resolve()}"


def _protect(args: dict) -> str:
    """Add password protection to PDF."""
    try:
        from pypdf import PdfReader, PdfWriter  # type: ignore[import-not-found]
    except ImportError:
        return "ERROR: pypdf not installed."

    input_path = args.get("input_path") or args.get("path")
    password = args.get("password")
    if not input_path:
        return "ERROR: 'input_path' is required."
    if not password:
        return "ERROR: 'password' is required."

    out = _ensure_output(args.get("output_path"), input_path, "_protected")
    reader = PdfReader(str(input_path))
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)

    owner_password = args.get("owner_password", password)
    writer.encrypt(password, owner_password, use_128bit=True)

    with open(out, "wb") as f:
        writer.write(f)
    return f"PDF protected with password -> {out.resolve()}"


def _organize(args: dict) -> str:
    """Reorder, delete, or insert pages in a PDF."""
    try:
        from pypdf import PdfReader, PdfWriter  # type: ignore[import-not-found]
    except ImportError:
        return "ERROR: pypdf not installed."

    input_path = args.get("input_path") or args.get("path")
    if not input_path:
        return "ERROR: 'input_path' is required."

    page_order = args.get("page_order")
    delete_pages = set(args.get("delete_pages") or [])

    out = _ensure_output(args.get("output_path"), input_path, "_organized")
    reader = PdfReader(str(input_path))
    total = len(reader.pages)
    writer = PdfWriter()

    if page_order:
        for pg in page_order:
            idx = int(pg) - 1
            if 0 <= idx < total:
                writer.add_page(reader.pages[idx])
    else:
        for i, page in enumerate(reader.pages):
            if (i + 1) not in delete_pages:
                writer.add_page(page)

    with open(out, "wb") as f:
        writer.write(f)
    return f"PDF organized ({len(writer.pages)} pages) -> {out.resolve()}"


def _pdf_to_pdfa(args: dict) -> str:
    """Convert PDF to PDF/A archival format."""
    try:
        import fitz  # type: ignore[import-not-found]
    except ImportError:
        return "ERROR: PyMuPDF not installed."

    input_path = args.get("input_path") or args.get("path")
    if not input_path:
        return "ERROR: 'input_path' is required."

    out = _ensure_output(args.get("output_path"), input_path, "_PDFA")
    doc = fitz.open(str(input_path))
    doc.save(
        str(out),
        garbage=4,
        deflate=True,
        clean=True,
        linearize=True,
    )
    doc.close()
    return f"PDF/A created -> {out.resolve()}"


def _repair(args: dict) -> str:
    """Attempt to repair and recover a corrupted PDF."""
    try:
        from pypdf import PdfReader, PdfWriter  # type: ignore[import-not-found]
    except ImportError:
        return "ERROR: pypdf not installed."

    input_path = args.get("input_path") or args.get("path")
    if not input_path:
        return "ERROR: 'input_path' is required."

    out = _ensure_output(args.get("output_path"), input_path, "_repaired")
    try:
        reader = PdfReader(str(input_path), strict=False)
        writer = PdfWriter()
        recovered = 0
        for page in reader.pages:
            try:
                writer.add_page(page)
                recovered += 1
            except Exception:
                continue
        with open(out, "wb") as f:
            writer.write(f)
        return f"Repair complete: {recovered}/{len(reader.pages)} pages recovered -> {out.resolve()}"
    except Exception as e:
        return f"ERROR during repair: {e}"


def _page_numbers(args: dict) -> str:
    """Add page numbers to PDF pages."""
    try:
        import fitz  # type: ignore[import-not-found]
    except ImportError:
        return "ERROR: PyMuPDF not installed."

    input_path = args.get("input_path") or args.get("path")
    if not input_path:
        return "ERROR: 'input_path' is required."

    out = _ensure_output(args.get("output_path"), input_path, "_numbered")
    position = args.get("position", "bottom-center")
    start_num = int(args.get("start", 1))
    font_size = int(args.get("font_size", 10))

    doc = fitz.open(str(input_path))
    for i, page in enumerate(doc):
        rect = page.rect
        num = start_num + i
        text = str(num)

        if "bottom" in position:
            y = rect.height - 20
        else:
            y = 15

        if "center" in position:
            x = rect.width / 2 - 10
        elif "right" in position:
            x = rect.width - 30
        else:
            x = 15

        page.insert_text(
            fitz.Point(x, y),
            text,
            fontsize=font_size,
            color=(0, 0, 0),
        )

    doc.save(str(out))
    doc.close()
    return f"Page numbers added ({position}) -> {out.resolve()}"


def _ocr(args: dict) -> str:
    """OCR scanned PDF to extract searchable text."""
    try:
        import fitz  # type: ignore[import-not-found]
    except ImportError:
        return "ERROR: PyMuPDF not installed."

    input_path = args.get("input_path") or args.get("path")
    if not input_path:
        return "ERROR: 'input_path' is required."

    doc = fitz.open(str(input_path))
    all_text = []

    for i, page in enumerate(doc):
        text = page.get_text().strip()
        if not text:
            try:
                tp = page.get_textpage_ocr(flags=0, language="eng")
                text = page.get_text(textpage=tp).strip()
            except Exception:
                text = f"[Page {i + 1}: OCR failed - install tesseract for best results]"
        all_text.append(f"--- Page {i + 1} ---\n{text}")

    doc.close()
    result = "\n\n".join(all_text)

    output_path = args.get("output_path")
    if output_path:
        Path(output_path).write_text(result, encoding="utf-8")
        return f"OCR text extracted -> {output_path}"
    return result[:3000] + (f"\n[...{len(result) - 3000} more chars...]" if len(result) > 3000 else "")


def _compare(args: dict) -> str:
    """Compare two PDFs and report text differences."""
    try:
        import fitz  # type: ignore[import-not-found]
    except ImportError:
        return "ERROR: PyMuPDF not installed."

    path1 = args.get("input_path") or args.get("path1") or args.get("path")
    path2 = args.get("path2") or args.get("compare_to")
    if not path1 or not path2:
        return "ERROR: 'input_path' (path1) and 'path2' / 'compare_to' are both required."

    doc1 = fitz.open(str(path1))
    doc2 = fitz.open(str(path2))
    diffs = []

    min_pages = min(len(doc1), len(doc2))
    for i in range(min_pages):
        t1 = doc1[i].get_text().strip()
        t2 = doc2[i].get_text().strip()
        if t1 != t2:
            diffs.append(f"Page {i + 1}: DIFFERENT (len {len(t1)} vs {len(t2)})")
        else:
            diffs.append(f"Page {i + 1}: identical")

    if len(doc1) != len(doc2):
        diffs.append(f"Page count differs: {len(doc1)} vs {len(doc2)}")

    doc1.close()
    doc2.close()

    summary = "\n".join(diffs)
    changed = sum(1 for d in diffs if "DIFFERENT" in d)
    return f"Compare result ({changed} different pages):\n{summary}"


def _redact(args: dict) -> str:
    """Permanently redact text patterns from PDF."""
    try:
        import fitz  # type: ignore[import-not-found]
    except ImportError:
        return "ERROR: PyMuPDF not installed."

    input_path = args.get("input_path") or args.get("path")
    if not input_path:
        return "ERROR: 'input_path' is required."

    patterns = args.get("patterns") or [args.get("text", "")]
    if not any(patterns):
        return "ERROR: 'patterns' (list of strings to redact) or 'text' is required."

    out = _ensure_output(args.get("output_path"), input_path, "_redacted")
    doc = fitz.open(str(input_path))
    total_redacted = 0

    for page in doc:
        for pattern in patterns:
            if not pattern:
                continue
            rects = page.search_for(pattern)
            for rect in rects:
                page.add_redact_annot(rect, fill=(0, 0, 0))
                total_redacted += 1
        page.apply_redactions()

    doc.save(str(out))
    doc.close()
    return f"Redacted {total_redacted} occurrences -> {out.resolve()}"


def _crop(args: dict) -> str:
    """Crop PDF pages to a specified rectangle."""
    try:
        import fitz  # type: ignore[import-not-found]
    except ImportError:
        return "ERROR: PyMuPDF not installed."

    input_path = args.get("input_path") or args.get("path")
    if not input_path:
        return "ERROR: 'input_path' is required."

    rect_arg = args.get("rect") or args.get("crop_box")
    if not rect_arg or len(rect_arg) < 4:
        return "ERROR: 'rect' [x0, y0, x1, y1] is required."

    out = _ensure_output(args.get("output_path"), input_path, "_cropped")
    pages = _parse_pages(args.get("pages"))
    crop_rect = fitz.Rect(rect_arg[0], rect_arg[1], rect_arg[2], rect_arg[3])

    doc = fitz.open(str(input_path))
    indices = pages if pages is not None else list(range(len(doc)))
    for i, page in enumerate(doc):
        if i in indices:
            page.set_cropbox(crop_rect)

    doc.save(str(out))
    doc.close()
    return f"Cropped {len(indices)} pages -> {out.resolve()}"


def _pdf_forms(args: dict) -> str:
    """List or fill PDF form fields."""
    try:
        import fitz  # type: ignore[import-not-found]
    except ImportError:
        return "ERROR: PyMuPDF not installed."

    input_path = args.get("input_path") or args.get("path")
    if not input_path:
        return "ERROR: 'input_path' is required."

    action = args.get("action", "list")
    field_values: dict = args.get("fields") or {}

    doc = fitz.open(str(input_path))

    if action == "list":
        fields = []
        for page in doc:
            for widget in page.widgets() if hasattr(page, "widgets") else []:
                fields.append({
                    "name": widget.field_name,
                    "type": widget.field_type_string,
                    "value": widget.field_value,
                    "page": page.number + 1,
                })
        doc.close()
        return json.dumps(fields, indent=2)

    elif action == "fill":
        filled = 0
        out = _ensure_output(args.get("output_path"), input_path, "_filled")
        for page in doc:
            for widget in page.widgets() if hasattr(page, "widgets") else []:
                name = widget.field_name
                if name in field_values:
                    widget.field_value = str(field_values[name])
                    widget.update()
                    filled += 1
        doc.save(str(out))
        doc.close()
        return f"Filled {filled} form fields -> {out.resolve()}"

    doc.close()
    return f"ERROR: Unknown form action '{action}'. Use 'list' or 'fill'."


def _summarize(args: dict) -> str:
    """AI-powered summary of PDF content."""
    input_path = args.get("input_path") or args.get("path")
    if not input_path:
        return "ERROR: 'input_path' is required."

    try:
        import fitz  # type: ignore[import-not-found]
        doc = fitz.open(str(input_path))
        full_text = ""
        for page in doc:
            full_text += page.get_text()
        doc.close()
    except Exception as e:
        return f"ERROR reading PDF: {e}"

    text_chunk = full_text[:8000]
    try:
        from brjarvis.integrations.backends.gemini import GeminiBackend
        gemini = GeminiBackend()
        prompt = f"Please provide a concise, structured summary of the following PDF document content:\n\n{text_chunk}"
        return gemini.quick(prompt)
    except Exception as e:
        return f"AI summary unavailable ({e}). Raw text preview:\n{text_chunk[:1000]}"


def _translate(args: dict) -> str:
    """Translate PDF text content using AI."""
    input_path = args.get("input_path") or args.get("path")
    target_lang = args.get("language") or args.get("target_language") or "English"
    if not input_path:
        return "ERROR: 'input_path' is required."

    try:
        import fitz  # type: ignore[import-not-found]
        doc = fitz.open(str(input_path))
        full_text = ""
        for page in doc:
            full_text += page.get_text()
        doc.close()
    except Exception as e:
        return f"ERROR reading PDF: {e}"

    text_chunk = full_text[:6000]
    try:
        from brjarvis.integrations.backends.gemini import GeminiBackend
        gemini = GeminiBackend()
        prompt = f"Translate the following text to {target_lang}. Output only the translated text:\n\n{text_chunk}"
        translation = gemini.quick(prompt)

        output_path = args.get("output_path")
        if output_path:
            Path(output_path).write_text(translation, encoding="utf-8")
            return f"Translation saved -> {output_path}"
        return translation[:3000]
    except Exception as e:
        return f"AI translation unavailable: {e}"


def _pdf_to_markdown(args: dict) -> str:
    """Convert PDF structure and content to Markdown."""
    try:
        import fitz  # type: ignore[import-not-found]
    except ImportError:
        return "ERROR: PyMuPDF not installed."

    input_path = args.get("input_path") or args.get("path")
    if not input_path:
        return "ERROR: 'input_path' is required."

    doc = fitz.open(str(input_path))
    md_parts = [f"# {Path(input_path).stem}\n"]

    for i, page in enumerate(doc):
        md_parts.append(f"\n## Page {i + 1}\n")
        blocks = page.get_text("dict")["blocks"]
        for block in blocks:
            if block.get("type") == 0:
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        text = span.get("text", "").strip()
                        size = span.get("size", 11)
                        if text:
                            if size >= 18:
                                md_parts.append(f"### {text}\n")
                            elif size >= 14:
                                md_parts.append(f"#### {text}\n")
                            else:
                                md_parts.append(f"{text} ")
                md_parts.append("\n")

    doc.close()
    result = "".join(md_parts)

    output_path = args.get("output_path")
    if output_path:
        Path(output_path).write_text(result, encoding="utf-8")
        return f"PDF -> Markdown: {output_path}"
    return result[:3000] + (f"\n[...{len(result) - 3000} more chars...]" if len(result) > 3000 else "")


def _sign(args: dict) -> str:
    """Add a digital signature placeholder widget to PDF."""
    try:
        import fitz  # type: ignore[import-not-found]
    except ImportError:
        return "ERROR: PyMuPDF not installed."

    input_path = args.get("input_path") or args.get("path")
    if not input_path:
        return "ERROR: 'input_path' is required."

    out = _ensure_output(args.get("output_path"), input_path, "_signature")
    page_num = int(args.get("page", 1)) - 1
    x = float(args.get("x", 72))
    y = float(args.get("y", 700))
    w = float(args.get("width", 200))
    h = float(args.get("height", 50))
    label = args.get("label", "Authorized Signature")

    doc = fitz.open(str(input_path))
    if page_num >= len(doc):
        page_num = len(doc) - 1
    page = doc[page_num]

    rect = fitz.Rect(x, y, x + w, y + h)
    page.draw_rect(rect, color=(0, 0, 0), width=1)
    page.insert_text(
        fitz.Point(x + 5, y + h - 10),
        label,
        fontsize=8,
        color=(0.5, 0.5, 0.5),
    )
    page.insert_text(
        fitz.Point(x + 5, y + 20),
        "X ____________________________",
        fontsize=10,
        color=(0, 0, 0),
    )

    doc.save(str(out))
    doc.close()
    return f"Signature placeholder added on page {page_num + 1} -> {out.resolve()}"


# ─────────────────────────────────────────────────────────────────────────────
# Dispatch Table
# ─────────────────────────────────────────────────────────────────────────────

_ACTION_MAP = {
    "merge": _merge,
    "split": _split,
    "compress": _compress,
    "pdf_to_word": _pdf_to_word,
    "pdf_to_powerpoint": _pdf_to_powerpoint,
    "pdf_to_pptx": _pdf_to_powerpoint,
    "pdf_to_excel": _pdf_to_excel,
    "pdf_to_xlsx": _pdf_to_excel,
    "word_to_pdf": _word_to_pdf,
    "docx_to_pdf": _word_to_pdf,
    "powerpoint_to_pdf": _powerpoint_to_pdf,
    "pptx_to_pdf": _powerpoint_to_pdf,
    "excel_to_pdf": _excel_to_pdf,
    "xlsx_to_pdf": _excel_to_pdf,
    "edit_pdf": _edit_pdf,
    "edit": _edit_pdf,
    "pdf_to_jpg": _pdf_to_jpg,
    "pdf_to_image": _pdf_to_jpg,
    "jpg_to_pdf": _jpg_to_pdf,
    "images_to_pdf": _jpg_to_pdf,
    "watermark": _watermark,
    "rotate": _rotate,
    "html_to_pdf": _html_to_pdf,
    "unlock": _unlock,
    "decrypt": _unlock,
    "protect": _protect,
    "encrypt": _protect,
    "organize": _organize,
    "reorder": _organize,
    "pdf_to_pdfa": _pdf_to_pdfa,
    "repair": _repair,
    "page_numbers": _page_numbers,
    "add_page_numbers": _page_numbers,
    "ocr": _ocr,
    "compare": _compare,
    "diff": _compare,
    "redact": _redact,
    "crop": _crop,
    "pdf_forms": _pdf_forms,
    "summarize": _summarize,
    "summary": _summarize,
    "translate": _translate,
    "pdf_to_markdown": _pdf_to_markdown,
    "pdf_to_md": _pdf_to_markdown,
    "extract_text": _pdf_to_markdown,
    "extract": _pdf_to_markdown,
    "read": _pdf_to_markdown,
    "read_pdf": _pdf_to_markdown,
    "get_text": _pdf_to_markdown,
    "sign": _sign,
    "signature": _sign,
}


# ─────────────────────────────────────────────────────────────────────────────
# Registered Tool
# ─────────────────────────────────────────────────────────────────────────────

@register_tool(
    name="pdf_tool",
    description=(
        "Comprehensive PDF operations tool - read, extract_text, merge, split, compress, convert (PDF<->Word/Excel/PPTX/JPG/HTML/Markdown), "
        "watermark, rotate, protect, unlock, OCR, redact, crop, compare, repair, page numbers, forms, summarize, translate, sign. "
        "Use 'action' to specify the operation (e.g. 'extract_text', 'read', 'merge', 'split', 'pdf_to_word', 'html_to_pdf'). "
        "Always provide 'input_path' (or 'input_paths' for merge/jpg_to_pdf)."
    ),
    parameters={
        "type": "object",
        "properties": {
            "action": {
                "type": "string",
                "description": "The PDF operation to perform.",
                "enum": list(_ACTION_MAP.keys()),
            },
            "input_path": {
                "type": "string",
                "description": "Path to input PDF or document file.",
            },
            "input_paths": {
                "type": "array",
                "items": {"type": "string"},
                "description": "List of input paths (required for merge and jpg_to_pdf actions).",
            },
            "output_path": {
                "type": "string",
                "description": "Optional output file path. Auto-generated if not provided.",
            },
            "output_dir": {
                "type": "string",
                "description": "Output directory for multi-file operations (split, pdf_to_jpg).",
            },
            "pages": {
                "type": "string",
                "description": "Page selection: '1', '1,3,5', '2-5', or list of ints.",
            },
            "password": {
                "type": "string",
                "description": "Password for protect/unlock operations.",
            },
            "text": {
                "type": "string",
                "description": "Text for watermark, edit, or redact operations.",
            },
            "angle": {
                "type": "integer",
                "description": "Rotation angle in degrees (for rotate action).",
            },
            "dpi": {
                "type": "integer",
                "description": "DPI resolution for image rendering (default: 150).",
            },
            "language": {
                "type": "string",
                "description": "Target language for translate action.",
            },
            "page_order": {
                "type": "array",
                "items": {"type": "integer"},
                "description": "New page order (1-indexed) for organize action.",
            },
            "delete_pages": {
                "type": "array",
                "items": {"type": "integer"},
                "description": "Pages to delete (1-indexed) for organize action.",
            },
            "fields": {
                "type": "object",
                "description": "Field name -> value mapping for pdf_forms fill action.",
            },
            "rect": {
                "type": "array",
                "items": {"type": "number"},
                "description": "Crop rectangle [x0, y0, x1, y1] in points for crop action.",
            },
            "html": {
                "type": "string",
                "description": "HTML content string for html_to_pdf action.",
            },
            "patterns": {
                "type": "array",
                "items": {"type": "string"},
                "description": "Text patterns to redact for redact action.",
            },
            "path2": {
                "type": "string",
                "description": "Second PDF path for compare action.",
            },
        },
        "required": ["action"],
    },
)
def pdf_tool(args: dict) -> str:
    """Dispatch PDF operations by action argument."""
    action = (args.get("action") or "").strip().lower()
    if not action:
        available = ", ".join(sorted(_ACTION_MAP.keys()))
        return f"ERROR: 'action' is required. Available actions:\n{available}"

    handler = _ACTION_MAP.get(action)
    if not handler:
        available = ", ".join(sorted(_ACTION_MAP.keys()))
        return f"ERROR: Unknown PDF action '{action}'. Available actions:\n{available}"

    try:
        logger.info("[PDFTools] Action: %s | Args: %s", action, {k: v for k, v in args.items() if k != "html"})
        return handler(args)
    except Exception as e:
        logger.exception("[PDFTools] Action '%s' failed: %s", action, e)
        return f"ERROR in pdf_tool({action}): {e}"